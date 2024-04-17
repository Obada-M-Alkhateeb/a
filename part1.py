# Import necessary libraries
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

# Function to extract tables from a PDF file
def extract_tables_from_pdf(pdf_path):
    return read_pdf(pdf_path, pages='all', multiple_tables=True)

# Function to add a table to a PowerPoint slide
def add_table_to_slide(slide, df):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(2), Inches(2), Inches(9), Inches(0.8)).table

    # Set column names as headers
    for col_index, column_name in enumerate(df.columns):
        cell = table.cell(0, col_index)
        cell.text = str(column_name)

    # Set row values
    for row_index, row in df.iterrows():
        for col_index, item in enumerate(row):
            cell = table.cell(row_index + 1, col_index)
            cell.text = str(item)

# Path to the PDF file
pdf_path = 'C:\\Users\\LENOVO\\Desktop\\New folder\\aaaa.pdf'

# Extract tables from the PDF
tables = extract_tables_from_pdf(pdf_path)

# Create a new PowerPoint presentation
prs = Presentation()

# Add a slide for each extracted table
for i, table_df in enumerate(tables):
    slide_layout = prs.slide_layouts[5]  # Use the title and content layout
    slide = prs.slides.add_slide(slide_layout)
    add_table_to_slide(slide, table_df)

# Save the PowerPoint file
prs.save('output.pptx')
